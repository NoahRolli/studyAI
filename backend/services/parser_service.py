# Parser-Service: Liest Text aus verschiedenen Dateiformaten
# Unterstützt PDF, Word, PowerPoint, Excel, Bilder (OCR), Markdown und Text
# PyMuPDF ist optional — wenn nicht installiert, werden PDFs mit Fallback gelesen

from docx import Document as DocxDocument  # python-docx — Library zum Word-Lesen
from pptx import Presentation  # python-pptx — Library zum PowerPoint-Lesen
from openpyxl import load_workbook  # openpyxl — Library zum Excel-Lesen
import pytesseract  # OCR — Text aus Bildern extrahieren
from PIL import Image  # Pillow — Bildverarbeitung
from pathlib import Path

# PyMuPDF optional importieren (Build-Probleme auf manchen Systemen)
try:
    import fitz
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False


# Unterstützte Dateiformate
SUPPORTED_FORMATS = {
    ".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".md",
    ".png", ".jpg", ".jpeg"
}


def parse_file(file_path: str) -> str:
    """
    Liest eine Datei und gibt den extrahierten Text zurück.
    Erkennt automatisch das Format anhand der Dateiendung.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_FORMATS:
        raise ValueError(f"Dateityp '{suffix}' wird nicht unterstützt. "
                         f"Unterstützt: {', '.join(sorted(SUPPORTED_FORMATS))}")

    if suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix == ".docx":
        return _parse_docx(path)
    elif suffix == ".pptx":
        return _parse_pptx(path)
    elif suffix == ".xlsx":
        return _parse_xlsx(path)
    elif suffix in (".png", ".jpg", ".jpeg"):
        return _parse_image(path)
    elif suffix == ".md":
        return _parse_markdown(path)
    elif suffix == ".txt":
        return _parse_txt(path)


def _parse_pdf(path: Path) -> str:
    """
    Extrahiert Text aus einer PDF-Datei.
    Nutzt PyMuPDF wenn verfügbar, sonst einfachen Textextrakt-Fallback.
    """
    if HAS_PYMUPDF:
        text = ""
        doc = fitz.open(str(path))
        for page in doc:
            text += page.get_text()
            text += "\n\n"
        doc.close()
        return text.strip()

    # Fallback: PDF als Binär lesen und Text extrahieren
    # Einfacher Extrakt — funktioniert für Text-PDFs, nicht für gescannte
    try:
        raw = path.read_bytes()
        text_parts = []
        # Suche nach Text-Streams im PDF
        i = 0
        while i < len(raw):
            # PDF Text-Objekte beginnen mit BT und enden mit ET
            bt = raw.find(b'BT', i)
            if bt == -1:
                break
            et = raw.find(b'ET', bt)
            if et == -1:
                break
            # Text zwischen Klammern extrahieren
            block = raw[bt:et]
            for match_start in range(len(block)):
                if block[match_start:match_start + 1] == b'(':
                    depth = 1
                    match_end = match_start + 1
                    while match_end < len(block) and depth > 0:
                        if block[match_end:match_end + 1] == b'(' and block[match_end - 1:match_end] != b'\\':
                            depth += 1
                        elif block[match_end:match_end + 1] == b')' and block[match_end - 1:match_end] != b'\\':
                            depth -= 1
                        match_end += 1
                    try:
                        text_parts.append(block[match_start + 1:match_end - 1].decode('latin-1'))
                    except Exception:
                        pass
            i = et + 2
        result = " ".join(text_parts)
        if result.strip():
            return result.strip()
        raise ValueError(
            "PDF-Text konnte nicht extrahiert werden. "
            "Installiere PyMuPDF für bessere PDF-Unterstützung: pip3 install PyMuPDF"
        )
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(
            f"PDF-Parsing fehlgeschlagen: {e}. "
            "Installiere PyMuPDF: pip3 install PyMuPDF"
        )


def _parse_docx(path: Path) -> str:
    """
    Extrahiert Text aus einer Word-Datei (.docx).
    Liest jeden Absatz einzeln aus.
    """
    doc = DocxDocument(str(path))
    paragraphs = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            paragraphs.append(paragraph.text)

    return "\n\n".join(paragraphs)


def _parse_pptx(path: Path) -> str:
    """
    Extrahiert Text aus einer PowerPoint-Datei (.pptx).
    Geht jede Folie und jede Textbox durch.
    """
    prs = Presentation(str(path))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = [f"--- Folie {slide_num} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_content.append(paragraph.text)

        slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


def _parse_xlsx(path: Path) -> str:
    """
    Extrahiert Text aus einer Excel-Datei (.xlsx).
    Liest jedes Tabellenblatt Zeile für Zeile.
    """
    wb = load_workbook(str(path), read_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_content = [f"--- Blatt: {sheet_name} ---"]

        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                sheet_content.append(" | ".join(cells))

        sheets_text.append("\n".join(sheet_content))

    wb.close()
    return "\n\n".join(sheets_text)


def _parse_image(path: Path) -> str:
    """
    Extrahiert Text aus einem Bild via OCR (Optical Character Recognition).
    Benötigt Tesseract als externes Programm.
    """
    try:
        image = Image.open(str(path))
        text = pytesseract.image_to_string(image, lang="deu+eng")
        return text.strip()
    except Exception as e:
        raise ValueError(f"OCR fehlgeschlagen: {e}. "
                         "Ist Tesseract installiert? (brew install tesseract)")


def _parse_markdown(path: Path) -> str:
    """
    Liest eine Markdown-Datei und gibt den Rohtext zurück.
    Struktur bleibt erhalten, da sie für AI nützlich ist.
    """
    raw = path.read_text(encoding="utf-8")
    return raw


def _parse_txt(path: Path) -> str:
    """
    Liest eine einfache Textdatei.
    Versucht verschiedene Encodings falls UTF-8 nicht klappt.
    """
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")